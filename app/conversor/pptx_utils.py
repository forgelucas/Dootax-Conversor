def extrair_texto_imagens_pptx(caminho_pptx):
    from pptx import Presentation
    prs = Presentation(caminho_pptx)
    imagens_por_slide = []

    for slide in prs.slides:
        imagens_slide = []
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                image = shape.image
                imagens_slide.append({
                    "dados": image.blob,
                    "ext": image.ext
                })
        imagens_por_slide.append(imagens_slide)
        
    return imagens_por_slide

def extrair_texto_formatado_do_pptx(caminho_pptx):
    from pptx import Presentation
    prs = Presentation(caminho_pptx)
    slides_formatados = []

    for idx, slide in enumerate(prs.slides, start=1):
        conteudo_slide = [{'tipo': 'titulo', 'texto': f"Slide {idx}:"}]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for par in shape.text_frame.paragraphs:
                par_dict = {'tipo': 'paragrafo', 'runs': [], 'bullet': par.level > 0}
                for run in par.runs:
                    par_dict['runs'].append({'texto': run.text, 'bold': run.font.bold})
                conteudo_slide.append(par_dict)
        slides_formatados.append(conteudo_slide)

    return slides_formatados
